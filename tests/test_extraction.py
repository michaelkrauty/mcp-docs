"""Tests for document extraction."""

import json
import tempfile
from pathlib import Path

import pytest

from mcp_docs.extraction import ContentExtractor, extract_content
from mcp_docs.extraction.notebook import extract_ipynb
from mcp_docs.extraction.text import (
    _csv_to_markdown_table,
    _read_text_with_encoding_fallback,
    extract_csv,
    extract_markdown,
    extract_rtf,
    extract_text,
)
from mcp_docs.models import DocumentType, ExtractionError


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


class TestTextExtraction:
    """Tests for plain text extraction."""

    def test_extract_text(self, temp_dir: Path) -> None:
        """Can extract content from text file."""
        text_file = temp_dir / "sample.txt"
        text_file.write_text("Hello, World! This is a test.")

        content = extract_text(text_file)

        assert content.text == "Hello, World! This is a test."
        assert content.word_count == 6
        assert content.title is None
        assert content.page_count is None

    def test_extract_text_utf8(self, temp_dir: Path) -> None:
        """Can handle UTF-8 content."""
        text_file = temp_dir / "unicode.txt"
        text_file.write_text("Héllo Wörld! 你好世界 🌍")

        content = extract_text(text_file)
        assert "Héllo Wörld" in content.text
        assert "你好世界" in content.text
        assert "🌍" in content.text

    def test_extract_empty_file(self, temp_dir: Path) -> None:
        """Can handle empty files."""
        text_file = temp_dir / "empty.txt"
        text_file.write_text("")

        content = extract_text(text_file)
        assert content.text == ""
        assert content.word_count == 0


class TestMarkdownExtraction:
    """Tests for markdown extraction."""

    def test_extract_markdown_with_title(self, temp_dir: Path) -> None:
        """Can extract title from H1 heading."""
        md_file = temp_dir / "doc.md"
        md_file.write_text("# My Document\n\nThis is the content.")

        content = extract_markdown(md_file)

        assert content.title == "My Document"
        assert "This is the content" in content.text

    def test_extract_markdown_no_title(self, temp_dir: Path) -> None:
        """Handles markdown without H1."""
        md_file = temp_dir / "no_title.md"
        md_file.write_text("## Section\n\nJust content here.")

        content = extract_markdown(md_file)

        assert content.title is None
        assert "Section" in content.text

    def test_extract_markdown_word_count(self, temp_dir: Path) -> None:
        """Counts words correctly in markdown."""
        md_file = temp_dir / "words.md"
        md_file.write_text("# Title\n\nOne two three four five.")

        content = extract_markdown(md_file)
        # "#", "Title", "One", "two", "three", "four", "five." = 7 words
        assert content.word_count == 7


class TestPptxExtraction:
    """Tests for PPTX extraction."""

    def test_extract_pptx(self, temp_dir: Path) -> None:
        """Can extract content from PPTX file."""
        from pptx import Presentation
        from pptx.util import Inches

        # Create a simple PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Add a text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        tf.text = "Hello from PowerPoint"

        pptx_file = temp_dir / "test.pptx"
        prs.save(pptx_file)

        # Extract
        from mcp_docs.extraction.office import extract_pptx

        content = extract_pptx(pptx_file)

        assert "Hello from PowerPoint" in content.text
        assert content.page_count == 1  # One slide
        assert content.word_count >= 3

    def test_extract_pptx_multiple_slides(self, temp_dir: Path) -> None:
        """Can extract from multiple slides."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Add two slides
        for i in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = f"Slide {i + 1} content"

        pptx_file = temp_dir / "multi.pptx"
        prs.save(pptx_file)

        from mcp_docs.extraction.office import extract_pptx

        content = extract_pptx(pptx_file)

        assert "Slide 1 content" in content.text
        assert "Slide 2 content" in content.text
        assert content.page_count == 2

    def test_extract_pptx_with_metadata(self, temp_dir: Path) -> None:
        """Extracts metadata from PPTX."""
        from pptx import Presentation

        prs = Presentation()
        prs.core_properties.title = "My Presentation"
        prs.core_properties.author = "Test Author"

        pptx_file = temp_dir / "meta.pptx"
        prs.save(pptx_file)

        from mcp_docs.extraction.office import extract_pptx

        content = extract_pptx(pptx_file)

        assert content.title == "My Presentation"
        assert content.metadata.get("author") == "Test Author"


class TestContentExtractor:
    """Tests for the ContentExtractor class."""

    def test_auto_detect_txt(self, temp_dir: Path) -> None:
        """Auto-detects .txt files."""
        txt_file = temp_dir / "file.txt"
        txt_file.write_text("Plain text content")

        extractor = ContentExtractor()
        content = extractor.extract(txt_file)

        assert content.text == "Plain text content"

    def test_auto_detect_md(self, temp_dir: Path) -> None:
        """Auto-detects .md files."""
        md_file = temp_dir / "file.md"
        md_file.write_text("# Markdown Title\n\nContent here.")

        extractor = ContentExtractor()
        content = extractor.extract(md_file)

        assert content.title == "Markdown Title"

    def test_auto_detect_pptx(self, temp_dir: Path) -> None:
        """Auto-detects .pptx files."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "PowerPoint content"

        pptx_file = temp_dir / "file.pptx"
        prs.save(pptx_file)

        extractor = ContentExtractor()
        content = extractor.extract(pptx_file)

        assert "PowerPoint content" in content.text

    def test_explicit_type(self, temp_dir: Path) -> None:
        """Can specify explicit document type."""
        # File with .md extension but treated as plain text
        md_file = temp_dir / "file.md"
        md_file.write_text("# Not A Title\n\nJust text.")

        extractor = ContentExtractor()
        content = extractor.extract(md_file, doc_type=DocumentType.TXT)

        # As plain text, title should be None (no extraction)
        assert content.title is None

    def test_file_not_found(self, temp_dir: Path) -> None:
        """Raises FileNotFoundError for missing files."""
        extractor = ContentExtractor()

        with pytest.raises(FileNotFoundError):
            extractor.extract(temp_dir / "nonexistent.txt")

    def test_can_extract(self) -> None:
        """can_extract returns correct results for supported types."""
        extractor = ContentExtractor()

        assert extractor.can_extract(DocumentType.PDF) is True
        assert extractor.can_extract(DocumentType.DOCX) is True
        assert extractor.can_extract(DocumentType.PPTX) is True
        assert extractor.can_extract(DocumentType.TXT) is True
        assert extractor.can_extract(DocumentType.MD) is True
        # DOC/PPT supported via MarkItDown (best-effort for legacy formats)
        assert extractor.can_extract(DocumentType.DOC) is True
        assert extractor.can_extract(DocumentType.PPT) is True


class TestExtractContentFunction:
    """Tests for the extract_content convenience function."""

    def test_extract_content_text(self, temp_dir: Path) -> None:
        """Convenience function works for text files."""
        txt_file = temp_dir / "test.txt"
        txt_file.write_text("Test content for extraction.")

        content = extract_content(txt_file)

        assert content.text == "Test content for extraction."
        assert content.word_count == 4

    def test_extract_content_markdown(self, temp_dir: Path) -> None:
        """Convenience function works for markdown files."""
        md_file = temp_dir / "test.md"
        md_file.write_text("# Doc Title\n\nBody text here.")

        content = extract_content(md_file)

        assert content.title == "Doc Title"


class TestNotebookExtraction:
    """Tests for Jupyter notebook (.ipynb) extraction."""

    @staticmethod
    def _write(path: Path, notebook: dict) -> None:
        path.write_text(json.dumps(notebook), encoding="utf-8")

    def test_extracts_markdown_prose_and_code(self, temp_dir: Path) -> None:
        """Markdown becomes prose, code becomes fenced blocks tagged with the language."""
        nb = {
            "cells": [
                {"cell_type": "markdown", "source": ["# Analysis\n", "\n", "Some **prose** here."]},
                {"cell_type": "code", "source": ["import pandas as pd\n", "df = load()"]},
            ],
            "metadata": {"language_info": {"name": "python"}},
            "nbformat": 4,
        }
        nb_file = temp_dir / "analysis.ipynb"
        self._write(nb_file, nb)

        content = extract_ipynb(nb_file)

        assert content.title == "Analysis"  # from first H1
        assert "Some **prose** here." in content.text
        assert "```python" in content.text
        assert "import pandas as pd" in content.text
        assert content.metadata["language"] == "python"
        assert content.metadata["cell_count"] == 2
        assert content.page_count is None

    def test_string_source_and_skips_outputs_and_raw_cells(self, temp_dir: Path) -> None:
        """source may be a bare string; cell outputs and raw cells are excluded."""
        nb = {
            "cells": [
                {
                    "cell_type": "code",
                    "source": "print('hi')",
                    "outputs": [{"output_type": "stream", "text": "OUTPUT_SHOULD_NOT_APPEAR"}],
                },
                {"cell_type": "raw", "source": ["RAW_SHOULD_NOT_APPEAR"]},
                {"cell_type": "markdown", "source": ""},  # empty -> skipped
            ],
            "metadata": {},
        }
        nb_file = temp_dir / "mixed.ipynb"
        self._write(nb_file, nb)

        content = extract_ipynb(nb_file)

        assert "print('hi')" in content.text  # bare-string source handled
        assert "OUTPUT_SHOULD_NOT_APPEAR" not in content.text  # outputs excluded
        assert "RAW_SHOULD_NOT_APPEAR" not in content.text  # raw cell excluded
        assert content.title is None
        assert "language" not in content.metadata  # no language metadata present

    def test_language_falls_back_to_kernelspec(self, temp_dir: Path) -> None:
        """When language_info is absent, the kernelspec language is used for fences."""
        nb = {
            "cells": [{"cell_type": "code", "source": "x <- 1"}],
            "metadata": {"kernelspec": {"language": "r"}},
        }
        nb_file = temp_dir / "rlang.ipynb"
        self._write(nb_file, nb)

        content = extract_ipynb(nb_file)
        assert "```r" in content.text

    @pytest.mark.parametrize(
        ("kernel_name", "expected"),
        [
            ("python3", "python"),
            ("python3.11", "python"),
            ("julia-1.9", "julia"),
            ("ir", "ir"),
        ],
    )
    def test_language_falls_back_to_kernel_name(
        self, temp_dir: Path, kernel_name: str, expected: str
    ) -> None:
        """With only a kernelspec name (no language_info/language), the language is
        the leading alphabetic run of the kernel name (so versioned names like
        "python3.11" and "julia-1.9" still map to "python"/"julia")."""
        nb = {
            "cells": [{"cell_type": "code", "source": "x = 1"}],
            "metadata": {"kernelspec": {"name": kernel_name}},
        }
        nb_file = temp_dir / "kn.ipynb"
        self._write(nb_file, nb)

        content = extract_ipynb(nb_file)
        assert f"```{expected}" in content.text
        assert content.metadata["language"] == expected

    def test_no_cells_yields_empty_content(self, temp_dir: Path) -> None:
        """A notebook with no cells extracts to empty text without error."""
        nb_file = temp_dir / "empty.ipynb"
        self._write(nb_file, {"metadata": {}, "nbformat": 4})

        content = extract_ipynb(nb_file)
        assert content.text == ""
        assert content.word_count == 0
        assert content.metadata["cell_count"] == 0

    def test_invalid_json_raises(self, temp_dir: Path) -> None:
        """A non-JSON file raises ExtractionError (caught by the processing queue)."""
        nb_file = temp_dir / "bad.ipynb"
        nb_file.write_text("{ not valid json", encoding="utf-8")
        with pytest.raises(ExtractionError):
            extract_ipynb(nb_file)

    def test_json_array_is_not_a_notebook(self, temp_dir: Path) -> None:
        """Valid JSON that is not an object raises ExtractionError."""
        nb_file = temp_dir / "array.ipynb"
        nb_file.write_text("[1, 2, 3]", encoding="utf-8")
        with pytest.raises(ExtractionError):
            extract_ipynb(nb_file)

    def test_from_extension_maps_ipynb(self) -> None:
        assert DocumentType.from_extension(".ipynb") == DocumentType.IPYNB
        assert DocumentType.from_extension("ipynb") == DocumentType.IPYNB

    def test_can_extract_ipynb(self) -> None:
        assert ContentExtractor().can_extract(DocumentType.IPYNB) is True

    def test_extract_dispatches_ipynb(self, temp_dir: Path) -> None:
        """ContentExtractor.extract() auto-detects .ipynb and routes to the notebook extractor."""
        nb = {
            "cells": [{"cell_type": "markdown", "source": "# Title\n\nbody text"}],
            "metadata": {},
        }
        nb_file = temp_dir / "dispatch.ipynb"
        self._write(nb_file, nb)

        content = ContentExtractor().extract(nb_file)
        assert content.title == "Title"
        assert "body text" in content.text

    def test_nbformat_v3_worksheets_layout(self, temp_dir: Path) -> None:
        """Legacy nbformat v3: cells nest under worksheets and code uses 'input'."""
        nb = {
            "nbformat": 3,
            "worksheets": [
                {
                    "cells": [
                        {
                            "cell_type": "markdown",
                            "source": ["# V3 Notebook\n", "\n", "legacy prose"],
                        },
                        {
                            "cell_type": "code",
                            "input": ["print('v3 code')"],
                            "outputs": [{"output_type": "stream", "text": "V3_OUTPUT_HIDDEN"}],
                        },
                    ]
                }
            ],
            "metadata": {"language_info": {"name": "python"}},
        }
        nb_file = temp_dir / "legacy.ipynb"
        self._write(nb_file, nb)

        content = extract_ipynb(nb_file)

        assert content.title == "V3 Notebook"
        assert "legacy prose" in content.text
        assert "print('v3 code')" in content.text  # read from the v3 'input' field
        assert "V3_OUTPUT_HIDDEN" not in content.text  # outputs still excluded
        assert content.metadata["cell_count"] == 2

    def test_pathologically_nested_json_raises(self, temp_dir: Path) -> None:
        """Deeply nested JSON (RecursionError) surfaces as ExtractionError, not a crash."""
        nb_file = temp_dir / "nested.ipynb"
        nb_file.write_text("[" * 4000 + "]" * 4000, encoding="utf-8")
        with pytest.raises(ExtractionError):
            extract_ipynb(nb_file)


class TestCsvMarkdownTable:
    """_csv_to_markdown_table rendering (the fallback renderer)."""

    def test_basic_table(self) -> None:
        out = _csv_to_markdown_table("name,age\nAlice,30\nBob,25\n")
        assert "| name | age |" in out
        assert "| --- | --- |" in out
        assert "| Alice | 30 |" in out

    def test_quoted_field_with_embedded_comma(self) -> None:
        out = _csv_to_markdown_table('item,note\n"Smith, Inc",hello\n')
        assert "| Smith, Inc | hello |" in out

    def test_pipe_in_cell_is_escaped(self) -> None:
        out = _csv_to_markdown_table("a,b\nx|y,z\n")
        assert "x\\|y" in out

    def test_ragged_rows_padded(self) -> None:
        out = _csv_to_markdown_table("a,b,c\n1,2\n")
        assert "| 1 | 2 |  |" in out

    def test_empty(self) -> None:
        assert _csv_to_markdown_table("") == ""

    def test_cr_only_line_endings(self) -> None:
        out = _csv_to_markdown_table("a,b\r1,2\r3,4\r")
        assert "| 1 | 2 |" in out
        assert "| 3 | 4 |" in out


class TestEncodingFallback:
    """_read_text_with_encoding_fallback decodes regardless of encoding/locale."""

    def test_latin1(self, temp_dir: Path) -> None:
        f = temp_dir / "l.txt"
        f.write_bytes("José,Montréal".encode("latin-1"))
        out = _read_text_with_encoding_fallback(f)
        assert "José" in out and "Montréal" in out

    def test_cp1252_smart_quotes_before_latin1(self, temp_dir: Path) -> None:
        # latin-1 would map these to C1 control chars; cp1252 must win first.
        f = temp_dir / "c.txt"
        f.write_bytes(b"\x93hi\x94 \x8050")
        out = _read_text_with_encoding_fallback(f)
        assert "“hi”" in out and "€50" in out

    def test_utf8_bom_stripped(self, temp_dir: Path) -> None:
        f = temp_dir / "b.txt"
        f.write_bytes("header".encode("utf-8-sig"))
        assert _read_text_with_encoding_fallback(f) == "header"

    def test_utf16_with_bom(self, temp_dir: Path) -> None:
        f = temp_dir / "u.txt"
        f.write_bytes("Zürich".encode("utf-16"))
        out = _read_text_with_encoding_fallback(f)
        assert out == "Zürich" and "\x00" not in out

    def test_bom_less_utf16le(self, temp_dir: Path) -> None:
        f = temp_dir / "u2.txt"
        f.write_bytes("Alice,Paris".encode("utf-16-le"))
        out = _read_text_with_encoding_fallback(f)
        assert "Alice" in out and "Paris" in out and "\x00" not in out

    def test_bom_less_utf16be_endian(self, temp_dir: Path) -> None:
        # The NUL-pattern path must pick the correct endian, not accept UTF-16BE
        # input mis-decoded as UTF-16LE mojibake.
        f = temp_dir / "u3.txt"
        f.write_bytes("Alice,Paris".encode("utf-16-be"))
        out = _read_text_with_encoding_fallback(f)
        assert "Alice" in out and "Paris" in out and "\x00" not in out


class TestCsvExtraction:
    """extract_csv renders a markdown table via the encoding-robust direct path.

    markitdown's CSV converter is unreliable for this corpus (raises on latin-1,
    mojibakes cp1252, leaves a UTF-8 BOM), so CSV extraction does not use it.
    """

    def test_returns_markdown_table(self, temp_dir: Path) -> None:
        f = temp_dir / "a.csv"
        f.write_text("name,age\nAlice,30\nBob,25\n")
        content = extract_csv(f)
        assert "| name | age |" in content.text
        assert "| Alice | 30 |" in content.text
        assert content.word_count > 0

    def test_latin1_extracts(self, temp_dir: Path) -> None:
        """The exact previously-failing case: accented names in a latin-1 export."""
        f = temp_dir / "latin1.csv"
        f.write_bytes("name,city\nJosé,Montréal\n".encode("latin-1"))
        content = extract_csv(f)
        assert "José" in content.text and "Montréal" in content.text

    def test_cp1252_decoded_cleanly(self, temp_dir: Path) -> None:
        """cp1252 smart quotes/euro decode correctly (markitdown mojibakes these)."""
        f = temp_dir / "cp.csv"
        f.write_bytes(b"item,price\n\x93widget\x94,\x8050\n")
        content = extract_csv(f)
        assert "“widget”" in content.text and "€50" in content.text

    def test_utf8_bom_header_clean(self, temp_dir: Path) -> None:
        """No stray BOM in the first header cell (markitdown leaves one)."""
        f = temp_dir / "bom.csv"
        f.write_bytes("name,city\nZoe,Paris\n".encode("utf-8-sig"))
        content = extract_csv(f)
        assert content.text.splitlines()[0].startswith("| name |")

    def test_empty_csv(self, temp_dir: Path) -> None:
        f = temp_dir / "e.csv"
        f.write_text("")
        content = extract_csv(f)
        assert content.word_count == 0


class TestRtfExtraction:
    """extract_rtf shares the encoding-fallback helper after the refactor."""

    def test_rtf_extracts_text(self, temp_dir: Path) -> None:
        f = temp_dir / "a.rtf"
        f.write_text(r"{\rtf1\ansi\deff0 Hello World}")
        content = extract_rtf(f)
        assert "Hello" in content.text

    def test_rtf_non_ascii_cp1252(self, temp_dir: Path) -> None:
        f = temp_dir / "b.rtf"
        # cp1252 smart-quote byte (0x93) must not raise
        f.write_bytes(b"{\\rtf1\\ansi Hello \x93World\x94}")
        content = extract_rtf(f)
        assert "Hello" in content.text
